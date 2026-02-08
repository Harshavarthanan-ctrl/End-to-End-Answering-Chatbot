import os
import pypdf
import docx
from pptx import Presentation

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            return extract_pdf(file_path)
        elif ext == ".docx":
            return extract_docx(file_path)
        elif ext == ".pptx":
            return extract_pptx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        else:
            return "Unsupported file format."
    except Exception as e:
        return f"Error reading file: {str(e)}"

def extract_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
